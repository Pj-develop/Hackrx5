# Import necessary libraries
import pytesseract
from PIL import Image
import os
import PyPDF2
from pptx import Presentation
import tkinter as tk
from tkinter import filedialog

# Initialize Tkinter for file dialog
root = tk.Tk()
root.withdraw()  # Hide the main window

# Prompt the user to select a file
file_path = filedialog.askopenfilename(
    title="Select a file",
    filetypes=[("PDF files", ".pdf"), ("Image files", ".jpg .jpeg *.png"), ("PPTX files", ".pptx")]
)

# Ensure the user selected a file
if not file_path:
    print("No file selected. Exiting...")
    exit()

# Function to extract text from JPG using OCR
def extract_text_from_jpg(jpg_file_name):
    image = Image.open(jpg_file_name)
    extracted_text = pytesseract.image_to_string(image)
    return extracted_text

# Function to extract text from PDF
def extract_text_from_pdf(pdf_file_name):
    with open(pdf_file_name, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        extracted_text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            extracted_text += page.extract_text()
    return extracted_text

# Function to extract text from PPTX
def extract_text_from_pptx(pptx_file_name):
    presentation = Presentation(pptx_file_name)
    extracted_text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"
    return extracted_text

# Determine the file type and extract text accordingly
if file_path.endswith('.pdf'):
    extracted_text = extract_text_from_pdf(file_path)
elif file_path.endswith('.jpg') or file_path.endswith('.jpeg') or file_path.endswith('.png'):
    extracted_text = extract_text_from_jpg(file_path)
elif file_path.endswith('.pptx'):
    extracted_text = extract_text_from_pptx(file_path)
else:
    extracted_text = "Unsupported file format."

# After extracting the text
output_file_path = "output.txt"

# Write the extracted text to a .txt file
with open(output_file_path, 'w', encoding='utf-8') as output_file:
    output_file.write(extracted_text)

print(f"Extracted text has been written to {output_file_path}")
