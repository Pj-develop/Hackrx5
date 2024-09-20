from flask import Flask, request, jsonify
import pytesseract
from PIL import Image
import os
import PyPDF2
from pptx import Presentation
from io import BytesIO
from sentence_transformers import SentenceTransformer
import numpy as np
import faiss
import pickle
import logging

app = Flask(__name__)
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize the sentence transformer model
model = SentenceTransformer('all-MiniLM-L6-v2')
VECTOR_SIZE = model.get_sentence_embedding_dimension()
DOCS_FILE = "documents.pkl"

class KnowledgeBase:
    def __init__(self):
        self.vector_store = VectorStore(VECTOR_SIZE)
        self.documents = self._load_or_create_documents()

    def _load_or_create_documents(self):
        if os.path.exists(DOCS_FILE):
            logger.info(f"Loading existing documents from {DOCS_FILE}")
            with open(DOCS_FILE, 'rb') as f:
                return pickle.load(f)
        else:
            logger.info("Creating new document storage")
            return []

    def _save_documents(self):
        with open(DOCS_FILE, 'wb') as f:
            pickle.dump(self.documents, f)
        logger.info(f"Documents saved to {DOCS_FILE}")

    def add_document(self, text, metadata):
        vector = model.encode(text)
        self.vector_store.add_vector(vector)
        self.documents.append({"text": text, **metadata})
        self._save_documents()

    def query(self, query_text, top_k=1):
        query_vector = model.encode(query_text)
        distances, indices = self.vector_store.search(query_vector, top_k)
        results = []
        for i, idx in enumerate(indices[0]):
            if idx != -1:
                results.append({
                    "score": 1 - distances[0][i],
                    "document": self.documents[idx]
                })
        return results

class VectorStore:
    def __init__(self, dimension, index_file='faiss.index'):
        self.dimension = dimension
        self.index_file = index_file
        self.index = self._load_or_create_index()

    def _load_or_create_index(self):
        if os.path.exists(self.index_file):
            return faiss.read_index(self.index_file)
        else:
            return faiss.IndexFlatL2(self.dimension)

    def add_vector(self, vector):
        self.index.add(np.array([vector]))
        self._save_index()

    def search(self, query_vector, k):
        return self.index.search(np.array([query_vector]), k)

    def _save_index(self):
        faiss.write_index(self.index, self.index_file)

kb = KnowledgeBase()

# Helper functions to extract text from files
def extract_text_from_jpg(jpg_file):
    image = Image.open(jpg_file)
    extracted_text = pytesseract.image_to_string(image)
    return extracted_text

def extract_text_from_pdf(pdf_file):
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    extracted_text = ""
    for page_num in range(len(pdf_reader.pages)):
        page = pdf_reader.pages[page_num]
        extracted_text += page.extract_text()
    return extracted_text

def extract_text_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    extracted_text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"
    return extracted_text

# Route to upload a file and extract text
@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    # Extract text based on file type
    if file.filename.endswith('.pdf'):
        extracted_text = extract_text_from_pdf(file)
    elif file.filename.endswith(('.jpg', '.jpeg', '.png')):
        extracted_text = extract_text_from_jpg(file)
    elif file.filename.endswith('.pptx'):
        extracted_text = extract_text_from_pptx(file)
    else:
        return jsonify({"error": "Unsupported file format"}), 400

    # Add the extracted text to the knowledge base
    kb.add_document(extracted_text, {"source": file.filename})

    return jsonify({"message": "Text extracted and added to knowledge base"}), 200

# Example route to query the knowledge base
@app.route('/query', methods=['POST'])
def query_knowledge_base():
    data = request.json
    query_text = data.get('query', '')
    results = kb.query(query_text, top_k=2)
    return jsonify(results)

if __name__ == "__main__":
    app.run(debug=True)
